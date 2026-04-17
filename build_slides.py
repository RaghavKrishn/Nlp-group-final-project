"""
build_slides.py
Generates PDEBench-Lang.pptx — upload to Google Drive to open in Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)   # header bar
BLUE   = RGBColor(0x29, 0x80, 0xB9)   # accent / highlights
TEAL   = RGBColor(0x17, 0x9A, 0x8A)   # secondary accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2C, 0x3E, 0x50)   # body text
LGRAY  = RGBColor(0xEC, 0xF0, 0xF1)   # light background strip
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
AMBER  = RGBColor(0xE6, 0x7E, 0x22)

BASE = '/Users/krivansemlani/Nlp-group-final-project'
IMG = {
    'dialect'  : os.path.join(BASE, 'dialect_comparison.png'),
    'heatmap'  : os.path.join(BASE, 'cross_dialect_heatmap.png'),
    'zeroshot' : os.path.join(BASE, 'zeroshot_rouge_all.png'),
    'zs_preds' : os.path.join(BASE, 'zeroshot_family_preds.png'),
}

W = Inches(13.33)   # slide width  (16:9)
H = Inches(7.5)     # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def header_bar(slide, title, subtitle=None):
    """Dark navy header bar with white title."""
    add_rect(slide, 0, 0, W, Inches(1.35), fill_rgb=NAVY)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75),
             font_size=Pt(32), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.45),
                 font_size=Pt(16), color=RGBColor(0xAD, 0xD8, 0xE6), italic=True)


def bullet_block(slide, items, l, t, w, h,
                 font_size=Pt(17), color=DARK, spacing=Inches(0.38),
                 bullet_char='▸', bullet_color=BLUE):
    """Render a list of (indent_level, text) tuples as bullet points."""
    y = t
    for level, text in items:
        indent = Inches(level * 0.35)
        bx = slide.shapes.add_textbox(l + indent, y, w - indent, Inches(0.38))
        tf = bx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # bullet character
        r0 = p.add_run()
        r0.text = bullet_char + '  '
        r0.font.color.rgb = bullet_color
        r0.font.size = font_size
        r0.font.bold = True
        r1 = p.add_run()
        r1.text = text
        r1.font.size = font_size
        r1.font.color.rgb = color
        y += spacing


def code_box(slide, code_text, l, t, w, h, font_size=Pt(12)):
    add_rect(slide, l, t, w, h, fill_rgb=RGBColor(0xF4, 0xF6, 0xF7),
             line_rgb=RGBColor(0xBD, 0xC3, 0xC7), line_width=Pt(1))
    txBox = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.08),
                                     w - Inches(0.2), h - Inches(0.15))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.name  = 'Courier New'
    run.font.size  = font_size
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)


def colored_pill(slide, text, l, t, w, h, bg, fg=WHITE, font_size=Pt(14)):
    add_rect(slide, l, t, w, h, fill_rgb=bg)
    add_text(slide, text, l, t + Inches(0.04), w, h,
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# gradient-ish background: full navy
add_rect(slide, 0, 0, W, H, fill_rgb=NAVY)
# lighter panel bottom half
add_rect(slide, 0, Inches(4.2), W, Inches(3.3), fill_rgb=RGBColor(0x15, 0x2B, 0x52))

# decorative accent line
add_rect(slide, Inches(0.55), Inches(3.9), Inches(6.5), Inches(0.06), fill_rgb=TEAL)

add_text(slide, 'PDEBench-Lang',
         Inches(0.55), Inches(1.1), Inches(12), Inches(1.2),
         font_size=Pt(52), bold=True, color=WHITE)

add_text(slide, 'Does Notation Format Shape Neural Reasoning About PDEs?',
         Inches(0.55), Inches(2.35), Inches(11), Inches(0.85),
         font_size=Pt(24), color=RGBColor(0xAD, 0xD8, 0xE6), italic=True)

add_text(slide, 'Natural Language  ·  LaTeX  ·  Prefix  ·  Postfix',
         Inches(0.55), Inches(3.15), Inches(10), Inches(0.55),
         font_size=Pt(18), color=RGBColor(0x85, 0xC1, 0xE9))

add_text(slide, 'NLP Final Project  —  2026',
         Inches(0.55), Inches(4.55), Inches(6), Inches(0.45),
         font_size=Pt(15), color=RGBColor(0x7F, 0x8C, 0x8D))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Question
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'The Central Question',
           'Notation is just formatting — or does it change what the model learns?')

add_text(slide,
         'PDEs can be written in radically different ways. Does the format a model is trained on '
         'affect how well it reasons about equation structure?',
         Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.9),
         font_size=Pt(18), color=DARK)

# Four boxes showing same equation in 4 dialects
dialects = [
    ('Natural Language', 'The second time derivative of u\nequals c² times the second\nspatial derivative of u.', NAVY),
    ('LaTeX',            r'u_{tt} = c^{2} \, u_{xx}', BLUE),
    ('Prefix',           '=(d(d(u,t),t),\n *(^(c,2), d(d(u,x),x)))', TEAL),
    ('Postfix',          'u t d t d\nc 2 ^ u x d x d *\n=', RGBColor(0x8E, 0x44, 0xAD)),
]

bw = Inches(2.9)
for i, (label, code, col) in enumerate(dialects):
    lx = Inches(0.35 + i * 3.15)
    # header pill
    colored_pill(slide, label, lx, Inches(2.6), bw, Inches(0.38), col)
    # code area
    add_rect(slide, lx, Inches(2.98), bw, Inches(2.1),
             fill_rgb=RGBColor(0xF8, 0xF9, 0xFA),
             line_rgb=col, line_width=Pt(1.5))
    add_text(slide, code,
             lx + Inches(0.1), Inches(3.05), bw - Inches(0.2), Inches(1.9),
             font_size=Pt(13), color=DARK)

add_text(slide, 'Same equation — Wave PDE (u_tt = c²u_xx) — four representations',
         Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.4),
         font_size=Pt(14), color=RGBColor(0x7F, 0x8C, 0x8D),
         italic=True, align=PP_ALIGN.CENTER)

add_text(slide, 'We train one T5-small model per dialect and measure: accuracy, cross-dialect transfer, and zero-shot generalisation.',
         Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.7),
         font_size=Pt(16), color=DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset Design
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Dataset Design',
           '10,000 training instances · 3 held-out families · 4 dialects each')

# Left column — training families
add_rect(slide, Inches(0.35), Inches(1.5), Inches(5.8), Inches(5.7),
         fill_rgb=RGBColor(0xEA, 0xF2, 0xFB), line_rgb=BLUE, line_width=Pt(1))
add_text(slide, 'Training Families  (5 × 2,000 = 10,000)',
         Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.4),
         font_size=Pt(15), bold=True, color=NAVY)

train_fams = [
    ('Heat',      'u_t = α u_xx',             'diffusion'),
    ('Wave',      'u_tt = c² u_xx',            'wave propagation'),
    ('Burgers',   'u_t + u·u_x = ν u_xx',     'nonlinear transport'),
    ('Laplace',   'u_xx + u_yy = 0',           'steady-state'),
    ('Advection', 'u_t + c·u_x = 0',          'pure transport'),
]
for i, (name, eq, desc) in enumerate(train_fams):
    y = Inches(2.1 + i * 0.88)
    colored_pill(slide, name, Inches(0.5), y, Inches(1.4), Inches(0.35), NAVY, font_size=Pt(12))
    add_text(slide, eq,   Inches(2.0), y, Inches(2.2), Inches(0.4), font_size=Pt(13), color=DARK)
    add_text(slide, desc, Inches(4.25), y, Inches(1.6), Inches(0.4), font_size=Pt(12),
             color=RGBColor(0x7F, 0x8C, 0x8D), italic=True)

# Right column — held-out families
add_rect(slide, Inches(6.5), Inches(1.5), Inches(6.45), Inches(5.7),
         fill_rgb=RGBColor(0xFD, 0xF2, 0xF8), line_rgb=RED, line_width=Pt(1))
add_text(slide, 'Held-Out Families  (zero-shot test only)',
         Inches(6.65), Inches(1.6), Inches(6.1), Inches(0.4),
         font_size=Pt(15), bold=True, color=RED)

held_fams = [
    ('KleinGordon',        'u_tt − u_xx + m²u = 0',            'wave + mass term'),
    ('ReactionDiffusion',  'u_t = D·u_xx + r·u·(1−u)',         'diffusion + logistic growth'),
    ('Beam',               'u_tt + b·u_xxxx = 0',               'wave + 4th-order space'),
]
for i, (name, eq, desc) in enumerate(held_fams):
    y = Inches(2.2 + i * 1.4)
    colored_pill(slide, name, Inches(6.65), y, Inches(2.1), Inches(0.35), RED, font_size=Pt(11))
    add_text(slide, eq,   Inches(8.85), y, Inches(2.5), Inches(0.4), font_size=Pt(13), color=DARK)
    add_text(slide, desc, Inches(6.65), y + Inches(0.42), Inches(4.1), Inches(0.4),
             font_size=Pt(12), color=RGBColor(0x7F, 0x8C, 0x8D), italic=True)

add_text(slide, '⚠  Models never see these during training — used only for zero-shot evaluation',
         Inches(6.65), Inches(5.7), Inches(6.1), Inches(0.45),
         font_size=Pt(13), color=RED, bold=True)

# Output format box at bottom
add_rect(slide, Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.4),
         fill_rgb=RGBColor(0xF4, 0xF6, 0xF7))
add_text(slide,
         'Output format per instance:  family: Heat  |  operators: exp, polynomial  |  reasoning: Contains first-order time derivative…',
         Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.38),
         font_size=Pt(12), color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Model Setup
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Model Setup',
           'Four independent T5-small models — one per dialect')

# Architecture diagram: input → T5-small → output
boxes = [
    (Inches(0.4),  Inches(2.8), Inches(2.4), Inches(1.8), NAVY,  WHITE, 'Input\n(one dialect)'),
    (Inches(3.5),  Inches(2.8), Inches(2.8), Inches(1.8), BLUE,  WHITE, 'T5-small\nFine-tuned\n60M params'),
    (Inches(7.1),  Inches(2.8), Inches(5.6), Inches(1.8), TEAL,  WHITE, 'Output\nfamily: Wave\noperators: sin, cos, polynomial\nreasoning: Second-order in time…'),
]
for lx, ty, bw2, bh, bg, fg, txt in boxes:
    add_rect(slide, lx, ty, bw2, bh, fill_rgb=bg)
    add_text(slide, txt, lx + Inches(0.1), ty + Inches(0.15),
             bw2 - Inches(0.2), bh - Inches(0.3),
             font_size=Pt(14), bold=True, color=fg)

# Arrows
for ax in [Inches(2.95), Inches(6.45)]:
    add_text(slide, '→', ax, Inches(3.4), Inches(0.5), Inches(0.5),
             font_size=Pt(28), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Config bullets
bullet_block(slide, [
    (0, 'Base model: google/t5-small  (60M parameters)'),
    (0, 'Training: 3 epochs · batch size 16 · learning rate 3×10⁻⁴'),
    (0, 'Split: 80% train / 10% val / 10% test  (same split reused across all 4 models)'),
    (0, 'Inference on CPU to avoid MPS beam-search issues on Apple Silicon'),
], Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5),
   font_size=Pt(16))

add_text(slide, 'Everything else is identical — only the training dialect changes.',
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         font_size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Experiment 1: In-Dialect Baseline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Experiment 1 — In-Dialect Baseline',
           'Train on dialect X, test on dialect X')

# Four big accuracy pills
for i, (dialect, acc, col) in enumerate([
    ('Natural', '100%', NAVY),
    ('LaTeX',   '100%', BLUE),
    ('Prefix',  '98.5%', TEAL),
    ('Postfix', '98.2%', RGBColor(0x8E, 0x44, 0xAD)),
]):
    lx = Inches(0.4 + i * 3.2)
    colored_pill(slide, dialect, lx, Inches(1.65), Inches(2.9), Inches(0.45), col, font_size=Pt(16))
    add_text(slide, acc, lx, Inches(2.2), Inches(2.9), Inches(0.9),
             font_size=Pt(42), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, 'family accuracy', lx, Inches(3.05), Inches(2.9), Inches(0.4),
             font_size=Pt(13), color=RGBColor(0x7F, 0x8C, 0x8D),
             align=PP_ALIGN.CENTER, italic=True)

# Separator
add_rect(slide, Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.04),
         fill_rgb=LGRAY)

# Insight
add_text(slide, 'What this tells us',
         Inches(0.5), Inches(3.75), Inches(12), Inches(0.4),
         font_size=Pt(18), bold=True, color=NAVY)

bullet_block(slide, [
    (0, 'All four models saturate at ~100% — T5-small easily learns each dialect.'),
    (0, 'Trash score = 0% across all dialects: when the family label is correct, reasoning is also correct.'),
    (0, 'ROUGE-L ≈ 44% despite 100% accuracy — the model outputs a different valid reasoning template, not the exact ground-truth sentence.'),
    (0, 'Problem: a 100% ceiling tells us nothing about which dialect is better.'),
    (0, 'Solution → test each model on dialects it has never seen.'),
], Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0), font_size=Pt(16))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Experiment 2: Cross-Dialect Heatmap
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Experiment 2 — Cross-Dialect Generalisation',
           'Train on one dialect, test on all four → 4×4 accuracy matrix')

slide.shapes.add_picture(IMG['heatmap'],
                         Inches(0.3), Inches(1.45),
                         Inches(12.7), Inches(5.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cross-Dialect: Worked Example
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Cross-Dialect: What the Numbers Mean',
           'A Heat equation — tested on a model trained on Postfix')

add_text(slide, 'Input given to the model  (LaTeX dialect — model was trained on Postfix):',
         Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
         font_size=Pt(15), color=DARK, bold=True)
code_box(slide,
         r'\frac{\partial}{\partial t} u(t,x) = 1.42 \frac{\partial^{2}}{\partial x^{2}} u(t,x)',
         Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.55), font_size=Pt(14))

# Two output columns: postfix model vs natural model
for col_i, (title, label, reasoning, correct, col) in enumerate([
    ('Postfix model (78% cross-dialect accuracy on LaTeX)',
     'Heat', 'A single time derivative balanced against a second spatial derivative scaled by 1.42 identifies this as diffusive transport.', True, GREEN),
    ('Natural model (20% cross-dialect accuracy on LaTeX)',
     'Wave', 'Contains a second-order time derivative and second-order spatial derivative, indicating wave propagation.', False, RED),
]):
    lx = Inches(0.4 + col_i * 6.45)
    add_rect(slide, lx, Inches(2.65), Inches(6.2), Inches(4.4),
             fill_rgb=RGBColor(0xF0, 0xFF, 0xF0) if correct else RGBColor(0xFF, 0xF0, 0xF0),
             line_rgb=col, line_width=Pt(1.5))
    add_text(slide, title, lx + Inches(0.15), Inches(2.75),
             Inches(5.9), Inches(0.5),
             font_size=Pt(13), bold=True, color=col)
    add_text(slide, f'Predicted family:  {label}',
             lx + Inches(0.15), Inches(3.3), Inches(5.9), Inches(0.4),
             font_size=Pt(15), bold=True, color=DARK)
    add_text(slide, f'Reasoning:  "{reasoning}"',
             lx + Inches(0.15), Inches(3.75), Inches(5.9), Inches(1.5),
             font_size=Pt(13), color=DARK, italic=True)
    tick = '✓  Correct' if correct else '✗  Wrong'
    add_text(slide, tick,
             lx + Inches(0.15), Inches(5.3), Inches(5.9), Inches(0.45),
             font_size=Pt(18), bold=True, color=col)

add_text(slide,
         'Postfix learned PDE structure well enough to transfer across notation. '
         'Natural language model overfits to surface phrasing.',
         Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
         font_size=Pt(14), color=NAVY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Cross-Dialect: Key Findings
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Cross-Dialect: Key Findings')

rows = [
    ('Natural',  '100%', '20%',  '20%',  '20%',  '20.0%', NAVY),
    ('LaTeX',    '0%',   '100%', '0%',   '0%',   '0.0%',  BLUE),
    ('Prefix',   '0%',   '38.7%','98.5%','68.9%','35.9%', TEAL),
    ('Postfix',  '5.7%', '78.1%','77.9%','98.2%','53.9%', RGBColor(0x8E, 0x44, 0xAD)),
]
headers = ['Trained on', 'Natural', 'LaTeX', 'Prefix', 'Postfix', 'Avg off-diag']
col_w = [Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.8)]
col_x = [Inches(0.4)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# Header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, cx, Inches(1.55), cw, Inches(0.42), fill_rgb=NAVY)
    add_text(slide, hdr, cx + Inches(0.05), Inches(1.6), cw - Inches(0.1), Inches(0.35),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data rows
for ri, (name, n, l, p, po, avg, col) in enumerate(rows):
    vals = [name, n, l, p, po, avg]
    is_diag = ri   # diagonal index within natural/latex/prefix/postfix columns
    for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
        bg = RGBColor(0xF0, 0xF8, 0xFF) if ri % 2 == 0 else WHITE
        if j == is_diag + 1:   # diagonal cell
            bg = RGBColor(0xFF, 0xD7, 0x00)
        if j == 5:   # avg column
            bg = RGBColor(0xEA, 0xF2, 0xFB)
        add_rect(slide, cx, Inches(1.97 + ri * 0.55), cw, Inches(0.52), fill_rgb=bg,
                 line_rgb=RGBColor(0xD0, 0xD3, 0xD4), line_width=Pt(0.5))
        fc = col if j == 0 else DARK
        add_text(slide, val, cx + Inches(0.05), Inches(2.02 + ri * 0.55),
                 cw - Inches(0.1), Inches(0.42),
                 font_size=Pt(14), bold=(j == 0 or j == 5), color=fc,
                 align=PP_ALIGN.CENTER)

add_text(slide, 'Gold cells = in-dialect (diagonal)',
         Inches(0.4), Inches(4.4), Inches(5), Inches(0.35),
         font_size=Pt(13), color=DARK, italic=True)

bullet_block(slide, [
    (0, 'Postfix is the best cross-dialect generaliser  (avg 53.9% off-diagonal)'),
    (0, 'LaTeX completely fails cross-dialect  (0.0%) — memorised LaTeX token patterns'),
    (0, 'Prefix → Postfix = 68.9% and vice versa 77.9% — symbolic dialects transfer to each other'),
    (0, 'Natural → anything = 20% (random) — surface language doesn\'t transfer symbolically'),
], Inches(0.4), Inches(4.85), Inches(12.5), Inches(2.4), font_size=Pt(16))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Experiment 3: Zero-Shot Bar Chart
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Experiment 3 — Zero-Shot Generalisation',
           'All 4 models tested on 3 PDE families never seen during training')

slide.shapes.add_picture(IMG['zeroshot'],
                         Inches(0.2), Inches(1.45),
                         Inches(13.0), Inches(4.5))

slide.shapes.add_picture(IMG['zs_preds'],
                         Inches(0.2), Inches(5.8),
                         Inches(13.0), Inches(1.55))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Zero-Shot: Worked Example (Beam)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Zero-Shot: The Beam Equation — A Stress Test',
           'u_tt + b·u_xxxx = 0  (4th-order spatial derivative — completely new structure)')

add_text(slide, 'Same equation, four dialects:',
         Inches(0.5), Inches(1.5), Inches(12), Inches(0.38),
         font_size=Pt(15), bold=True, color=DARK)

inputs_and_results = [
    ('Natural',
     'u_tt plus 0.71 times the fourth derivative of u w.r.t. x equals zero.',
     'Advection', '✗', RED,
     '"Only first-order derivatives appear: u_t and 0.71*u_x…"',
     '"fourth" never seen in training → ignored'),
    ('LaTeX',
     r'\partial^2/\partial t^2 u + 0.71 \partial^4/\partial x^4 u = 0',
     'Wave', '✓~', AMBER,
     '"Contains second-order time derivative and second-order spatial derivative…"',
     'Sees \\partial^2 in time → calls it Wave (close enough)'),
    ('Prefix',
     '=(+(*(0.71, d(d(d(d(u,x),x),x),x)), d(d(u,t),t)), 0)',
     'Wave', '✓', GREEN,
     '"Contains second-order time derivative…"',
     'Counts 4 nested d() calls in x → recognises depth'),
    ('Postfix',
     '0.71 u x d x d x d x d * u t d t d + 0 =',
     'Advection', '✗', RED,
     '"Only first-order derivatives appear: u_t and 0.71*u_x…"',
     'Token repetition misread as first-order'),
]

for i, (dialect, inp, pred, tick, col, reasoning, why) in enumerate(inputs_and_results):
    lx = Inches(0.3 + (i % 2) * 6.5)
    ty = Inches(2.05 + (i // 2) * 2.6)
    bw2 = Inches(6.2)
    add_rect(slide, lx, ty, bw2, Inches(2.4),
             fill_rgb=RGBColor(0xF8, 0xF9, 0xFA),
             line_rgb=col, line_width=Pt(2))
    colored_pill(slide, f'{dialect}  {tick}', lx + Inches(0.1), ty + Inches(0.08),
                 Inches(2.0), Inches(0.33), col, font_size=Pt(12))
    add_text(slide, inp,
             lx + Inches(0.1), ty + Inches(0.48), bw2 - Inches(0.2), Inches(0.55),
             font_size=Pt(11), color=DARK)
    add_text(slide, f'→ Predicts: {pred}',
             lx + Inches(0.1), ty + Inches(1.05), bw2 - Inches(0.2), Inches(0.35),
             font_size=Pt(13), bold=True, color=col)
    add_text(slide, why,
             lx + Inches(0.1), ty + Inches(1.45), bw2 - Inches(0.2), Inches(0.8),
             font_size=Pt(12), color=RGBColor(0x5D, 0x6D, 0x7E), italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — The Big Reversal
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'The Big Reversal',
           'Postfix wins cross-dialect. Prefix wins zero-shot. Why?')

# Two big panels
for col_i, (title, subtitle, winner, loser, winner_col, loser_col,
            w_score, l_score, explanation) in enumerate([
    ('Cross-Dialect\nGeneralisation',
     'Same equation, different notation',
     'Postfix', 'LaTeX',
     RGBColor(0x8E, 0x44, 0xAD), RED,
     '53.9%', '0.0%',
     'Postfix encodes operator order as a left-to-right token sequence.\n'
     'This surface ordering is easy to re-match across notation styles.\n'
     'It generalises because it learned to recognise token patterns\n'
     'that are consistent across different dialects.'),
    ('Zero-Shot\nGeneralisation',
     'Never-seen PDE families',
     'Prefix', 'Postfix',
     TEAL, RED,
     '3/3 families', '1/3 families',
     'Prefix encodes operator depth as explicit nesting:\n'
     'd(d(u,x),x) = 2nd order.  d(d(d(d(u,x)…),x),x) = 4th order.\n'
     'The model learns to count nesting depth — which is exactly\n'
     'what determines PDE structural family.'),
]):
    lx = Inches(0.35 + col_i * 6.5)
    add_rect(slide, lx, Inches(1.5), Inches(6.2), Inches(5.65),
             fill_rgb=RGBColor(0xF8, 0xF9, 0xFA),
             line_rgb=winner_col, line_width=Pt(2))
    add_text(slide, title, lx + Inches(0.15), Inches(1.6),
             Inches(5.9), Inches(0.85),
             font_size=Pt(22), bold=True, color=NAVY)
    add_text(slide, subtitle, lx + Inches(0.15), Inches(2.45),
             Inches(5.9), Inches(0.35),
             font_size=Pt(14), color=RGBColor(0x7F, 0x8C, 0x8D), italic=True)

    colored_pill(slide, f'Best: {winner}  {w_score}',
                 lx + Inches(0.15), Inches(2.9), Inches(3.0), Inches(0.38),
                 winner_col, font_size=Pt(13))
    colored_pill(slide, f'Worst: {loser}  {l_score}',
                 lx + Inches(3.3), Inches(2.9), Inches(2.7), Inches(0.38),
                 loser_col, font_size=Pt(13))

    add_text(slide, explanation,
             lx + Inches(0.15), Inches(3.45), Inches(5.9), Inches(2.2),
             font_size=Pt(14), color=DARK)

add_text(slide,
         'Different notations encode different inductive biases.\n'
         'There is no single "best" dialect — the right choice depends on what kind of generalisation you need.',
         Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.6),
         font_size=Pt(15), bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Key Takeaways')

takeaways = [
    (NAVY,  '1', 'Notation format matters.',
     'All four models hit 100% in-dialect. The ceiling hides the real differences. '
     'Cross-dialect and zero-shot tests reveal that what the model actually learned is very different depending on dialect.'),
    (BLUE,  '2', 'Postfix → best cross-dialect generaliser  (53.9% avg off-diagonal).',
     'Symbolic token order is robust across notation styles. '
     'Prefix transfers to postfix and back. Natural and LaTeX each create isolated "notation islands".'),
    (TEAL,  '3', 'Prefix → best zero-shot generaliser  (3/3 unseen families correctly mapped).',
     'Nested function calls encode derivative order explicitly. '
     'Postfix fails on Beam and KleinGordon because repeated tokens look first-order. '
     'Structure-encoding beats pattern-matching for unseen equations.'),
]

for i, (col, num, title, body) in enumerate(takeaways):
    ty = Inches(1.7 + i * 1.7)
    colored_pill(slide, num, Inches(0.4), ty, Inches(0.55), Inches(0.9), col)
    add_text(slide, title,
             Inches(1.1), ty, Inches(11.5), Inches(0.5),
             font_size=Pt(20), bold=True, color=col)
    add_text(slide, body,
             Inches(1.1), ty + Inches(0.45), Inches(11.5), Inches(0.8),
             font_size=Pt(16), color=DARK)
    add_rect(slide, Inches(0.4), ty + Inches(1.35), Inches(12.5), Inches(0.03),
             fill_rgb=LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Future Work
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_rgb=WHITE)
header_bar(slide, 'Future Work & Limitations')

future = [
    ('Break the 100% ceiling',
     'Add ReactionDiffusion to training — gives the model a 6th family to distinguish, '
     'making in-dialect accuracy non-trivial and revealing dialect-specific weaknesses.'),
    ('Bigger model',
     'Repeat with T5-base (250M) or BART-base. Does scale change the cross-dialect gap? '
     'Does a larger model close the postfix zero-shot failure on Beam?'),
    ('Multiple seeds',
     'Run each training 3× with different random seeds to add error bars. '
     'Confirm that the prefix vs postfix zero-shot gap is statistically significant.'),
    ('Intermediate representations',
     'Train on a mix of two dialects (e.g., 50% prefix + 50% postfix). '
     'Does exposure to both symbolic forms further improve generalisation?'),
]

for i, (title, body) in enumerate(future):
    lx = Inches(0.4 + (i % 2) * 6.45)
    ty = Inches(1.65 + (i // 2) * 2.5)
    add_rect(slide, lx, ty, Inches(6.1), Inches(2.2),
             fill_rgb=RGBColor(0xEA, 0xF2, 0xFB),
             line_rgb=BLUE, line_width=Pt(1))
    add_text(slide, title,
             lx + Inches(0.15), ty + Inches(0.12), Inches(5.8), Inches(0.45),
             font_size=Pt(16), bold=True, color=NAVY)
    add_text(slide, body,
             lx + Inches(0.15), ty + Inches(0.6), Inches(5.8), Inches(1.4),
             font_size=Pt(14), color=DARK)

add_text(slide, 'Limitation: T5-small may not be powerful enough to fully exploit structural cues — '
         'a stronger baseline model is needed before drawing firm conclusions.',
         Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.6),
         font_size=Pt(14), color=RED, italic=True, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, 'PDEBench_Lang.pptx')
prs.save(out)
print(f'Saved: {out}')
